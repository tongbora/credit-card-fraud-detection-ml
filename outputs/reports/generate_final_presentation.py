from pathlib import Path
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.ticker import PercentFormatter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parents[2]
FIG_DIR = ROOT / "outputs" / "figures"
METRICS_FILE = ROOT / "outputs" / "metrics" / "model_comparison.csv"
REPORT_DIR = ROOT / "outputs" / "reports"
CHART_DIR = REPORT_DIR / "charts"
OUT_PPTX = REPORT_DIR / "Credit_Card_Fraud_Detection_Final_Presentation.pptx"


MODEL_COLORS = {
    "Logistic Regression": "#2563eb",
    "Random Forest": "#10b981",
    "XGBoost": "#f59e0b",
}


def safe_percent(v):
    return f"{float(v)*100:.2f}%"


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def style_title(shape):
    p = shape.text_frame.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)


def style_subtitle(shape):
    tf = shape.text_frame
    for p in tf.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(71, 85, 105)


def add_bullets(slide, title, bullets, notes=None):
    layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = title
    style_title(s.shapes.title)

    tf = s.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(21)
        p.font.color.rgb = RGBColor(30, 41, 59)
    if notes:
        add_notes(s, notes)
    return s


def add_title_slide():
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Credit Card Fraud Detection"
    style_title(s.shapes.title)
    sub = s.placeholders[1]
    sub.text = "Machine Learning Final Project\nStudent: ____________________\nCourse: ____________________"
    style_subtitle(sub)
    add_notes(
        s,
        "Introduce yourself and project context. Mention this is an end-to-end pipeline from training to dashboard deployment.",
    )


def add_image_slide(title, img_path, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    style_title(s.shapes.title)
    if img_path and img_path.exists():
        s.shapes.add_picture(str(img_path), Inches(0.8), Inches(1.4), width=Inches(12.0), height=Inches(5.7))
    else:
        box = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1))
        box.text_frame.text = "Image not found"
    if notes:
        add_notes(s, notes)
    return s


def add_two_image_slide(title, left_img, right_img, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    style_title(s.shapes.title)
    if left_img and left_img.exists():
        s.shapes.add_picture(str(left_img), Inches(0.6), Inches(1.45), width=Inches(6.1), height=Inches(5.5))
    if right_img and right_img.exists():
        s.shapes.add_picture(str(right_img), Inches(6.7), Inches(1.45), width=Inches(6.1), height=Inches(5.5))
    if notes:
        add_notes(s, notes)
    return s


def add_metrics_table_slide(metrics_df):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Comparative Analysis (Real Metrics)"
    style_title(s.shapes.title)

    rows = 4
    cols = 4
    table = s.shapes.add_table(rows, cols, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.7)).table

    headers = ["Model", "Precision", "Recall", "F1-score"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h

    models = ["Logistic Regression", "Random Forest", "XGBoost"]
    for r, model in enumerate(models, start=1):
        row = metrics_df.loc[model]
        table.cell(r, 0).text = model
        table.cell(r, 1).text = f"{row['precision']:.4f}"
        table.cell(r, 2).text = f"{row['recall']:.4f}"
        table.cell(r, 3).text = f"{row['f1']:.4f}"

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(16 if r == 0 else 15)
                p.font.bold = r == 0
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

    best_model = metrics_df['f1'].idxmax()
    best_f1 = metrics_df['f1'].max()
    text = s.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.8), Inches(1.2)).text_frame
    text.text = f"Best Model: {best_model}  |  Highest F1-score: {best_f1:.4f}"
    p = text.paragraphs[0]
    p.font.name = "Calibri"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(6, 95, 70)

    add_notes(
        s,
        "Use this as the main evidence slide. Highlight that Random Forest is selected from actual outputs/metrics/model_comparison.csv by highest F1-score.",
    )


def make_bar_chart(metrics_df, metric, title, out_path):
    plt.figure(figsize=(6.5, 4), dpi=180)
    models = metrics_df.index.tolist()
    vals = metrics_df[metric].values
    colors = [MODEL_COLORS.get(m, "#64748b") for m in models]
    best_idx = vals.argmax()
    colors[best_idx] = "#14b8a6"

    bars = plt.bar(models, vals, color=colors, edgecolor="#dbe4f0")
    plt.ylim(0, 1.05)
    plt.gca().yaxis.set_major_formatter(PercentFormatter(1.0))
    plt.title(title, fontsize=13, fontweight="bold")
    plt.grid(axis="y", color="#e5e7eb", linewidth=0.8)
    for b, v in zip(bars, vals):
        plt.text(b.get_x() + b.get_width()/2, v + 0.02, f"{v:.1%}", ha="center", fontsize=9, weight="bold")
    plt.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(out_path, bbox_inches="tight")
    plt.close()


def add_demo_slide(notes=None):
    s = add_bullets(
        None,
        "System Demo (Gradio Dashboard)",
        [
            "Model comparison dashboard with metric cards",
            "Single transaction prediction with fraud probability",
            "Batch CSV upload with fraud/normal summary",
            "Downloadable prediction results for reporting",
        ],
        notes=notes,
    )
    return s


# Build deck
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

metrics = pd.read_csv(METRICS_FILE, index_col=0)

add_title_slide()

add_bullets(
    None,
    "Introduction: Problem & Objective",
    [
        "Credit card fraud creates major financial and trust risks",
        "Need to catch fraudulent transactions early",
        "Need to reduce false alerts for legitimate users",
        "Objective: build a practical end-to-end ML fraud detection system",
    ],
    notes="Set business context and explain precision-recall tradeoff importance in fraud detection.",
)

add_bullets(
    None,
    "Dataset Overview",
    [
        "Source: Credit Card Fraud Detection dataset (Kaggle/ULB)",
        "284,807 transactions total",
        "492 fraud cases (0.17%)",
        "Features used: V1–V28 + Amount | Target: Class",
    ],
    notes="Emphasize severe imbalance. Mention why this invalidates naive accuracy-based conclusions.",
)

add_image_slide(
    "Data Insights: Class Imbalance & EDA",
    FIG_DIR / "01_class_distribution.png",
    notes="Use this chart to motivate imbalance handling with SMOTE.",
)

add_two_image_slide(
    "Data Insights: Amount & Correlations",
    FIG_DIR / "02_amount_distribution.png",
    FIG_DIR / "03_top_correlations.png",
    notes="Summarize EDA findings only from real figures: skewed amounts and multi-feature fraud signal.",
)

add_bullets(
    None,
    "Methodology: Preprocessing & SMOTE",
    [
        "Data quality checks (duplicates/missing values)",
        "Feature scaling with StandardScaler",
        "Stratified train-test split for class preservation",
        "SMOTE applied to training data to address imbalance",
        "Leakage-safe evaluation on untouched test set",
    ],
    notes="Explain why SMOTE is applied to train only and why leakage prevention matters.",
)

add_bullets(
    None,
    "Models Trained",
    [
        "Logistic Regression (baseline linear model)",
        "Random Forest (ensemble trees)",
        "XGBoost (gradient boosting)",
        "Randomized hyperparameter tuning in training pipeline",
    ],
    notes="Briefly explain rationale for each model family and practical runtime tuning strategy.",
)

add_metrics_table_slide(metrics)

# Generate metric charts
prec_chart = CHART_DIR / "precision_comparison.png"
rec_chart = CHART_DIR / "recall_comparison.png"
f1_chart = CHART_DIR / "f1_comparison.png"
make_bar_chart(metrics, "precision", "Precision Comparison", prec_chart)
make_bar_chart(metrics, "recall", "Recall Comparison", rec_chart)
make_bar_chart(metrics, "f1", "F1-score Comparison", f1_chart)

add_three = prs.slides.add_slide(prs.slide_layouts[5])
add_three.shapes.title.text = "Visual Analytics: Metric Comparison"
style_title(add_three.shapes.title)
add_three.shapes.add_picture(str(prec_chart), Inches(0.4), Inches(1.45), width=Inches(4.2), height=Inches(5.4))
add_three.shapes.add_picture(str(rec_chart), Inches(4.6), Inches(1.45), width=Inches(4.2), height=Inches(5.4))
add_three.shapes.add_picture(str(f1_chart), Inches(8.8), Inches(1.45), width=Inches(4.2), height=Inches(5.4))
add_notes(add_three, "Walk chart-by-chart and conclude Random Forest has strongest F1 balance in this run.")

add_two_image_slide(
    "ROC & PR Curves",
    FIG_DIR / "roc_curve_comparison.png",
    FIG_DIR / "pr_curve_comparison.png",
    notes="Explain ROC and PR relevance for imbalanced classification.",
)

add_demo_slide(
    notes="Live demo flow: model comparison section, single prediction, batch prediction + downloadable output."
)

best_model = metrics["f1"].idxmax()
best_row = metrics.loc[best_model]
add_bullets(
    None,
    "Conclusion",
    [
        "End-to-end fraud detection system built and deployed with Gradio",
        "Severe class imbalance handled using SMOTE",
        f"Best model by F1-score: {best_model} ({best_row['f1']:.4f})",
        f"Random Forest precision: {best_row['precision']:.4f} | recall: {best_row['recall']:.4f}",
    ],
    notes="State final recommendation clearly and reference real metrics only.",
)

add_bullets(
    None,
    "Future Improvements",
    [
        "Threshold optimization based on business cost",
        "Model explainability with SHAP/LIME",
        "Probability calibration and ensemble stacking",
        "Production monitoring, drift detection, and retraining",
    ],
    notes="Show roadmap from academic project to production maturity.",
)

add_bullets(
    None,
    "Q&A",
    [
        "Why not rely on accuracy?",
        "Why SMOTE and how leakage was prevented?",
        "Why Random Forest selected in final run?",
        "How would you deploy this in production?",
    ],
    notes="Keep answers concise and tied to your actual pipeline outputs.",
)

REPORT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPTX)
print(f"Saved PPTX: {OUT_PPTX}")
print(f"Generated charts: {prec_chart}, {rec_chart}, {f1_chart}")
